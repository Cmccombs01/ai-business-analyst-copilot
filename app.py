import streamlit as st
import os
import pandas as pd
import requests
import sqlite3
import subprocess
import yfinance as yf
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from groq import Groq
from fpdf import FPDF
from pptx import Presentation
from gtts import gTTS 
import base64
from audio_recorder_streamlit import audio_recorder

# 1. Setup & Authentication
load_dotenv()
my_api_key = os.getenv("GROQ_API_KEY") 
client = Groq(api_key=my_api_key)

# 2. Page Configuration & Custom CSS Cleanup
st.set_page_config(page_title="AI Business Analyst", page_icon="📈", layout="wide")

# HIDE STREAMLIT WATERMARKS & CLEAN UP PADDING
st.markdown("""
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    .block-container {padding-top: 2rem; padding-bottom: 2rem;}
    </style>
    """, unsafe_allow_html=True)

st.title("📈 AI Business Analyst Copilot")
st.caption("Powered by Llama 3.3, Whisper AI, and Python")

# 3. Sleek Sidebar Redesign
audio_prompt = None
with st.sidebar:
    st.header("⚙️ Copilot Settings")
    enable_voice = st.toggle("🔊 Enable AI Voice Responses", value=True)
    
    st.divider()
    
    st.write("🎙️ **Voice Interface**")
    audio_bytes = audio_recorder(text="Click mic to talk...", icon_size="2x")
    
    if audio_bytes:
        with st.spinner("Listening..."):
            with open("temp_voice.wav", "wb") as f:
                f.write(audio_bytes)
            try:
                with open("temp_voice.wav", "rb") as f:
                    transcription = client.audio.transcriptions.create(
                        file=("temp_voice.wav", f.read()),
                        model="whisper-large-v3"
                    )
                audio_prompt = transcription.text
            except Exception as e:
                st.error(f"Could not hear you properly: {e}")

    st.divider()
    
    # Cleaned up Command List using an Expander
    with st.expander("🛠️ View Available Commands"):
        st.markdown("""
        * `analyze csv [file.csv]`
        * `plot [file.csv]`
        * `predict [file.csv]`
        * `sentiment [file.csv]`
        * `report [file.csv]`
        * `deck [file.csv]`
        * `stock [ticker]`
        * `query db [question]`
        * `scrape [URL]`
        """)

# 4. Initialize AI Memory
sys_prompt = "You are a Senior Business Analyst mentoring Caleb. Keep answers professional and concise. Output code clearly when asked."
if "conversation_history" not in st.session_state:
    st.session_state.conversation_history = [{"role": "system", "content": sys_prompt}]
if "chat_display" not in st.session_state:
    st.session_state.chat_display = []

# Display previous chat messages with AVATARS
for msg in st.session_state.chat_display:
    avatar_icon = "👤" if msg["role"] == "user" else "🤖"
    with st.chat_message(msg["role"], avatar=avatar_icon):
        st.markdown(msg["content"])

# 5. The Chat Input Box
typed_prompt = st.chat_input("Ask your Copilot something...")
final_prompt = audio_prompt or typed_prompt

if final_prompt:
    # Display User Message
    with st.chat_message("user", avatar="👤"):
        st.markdown(final_prompt)
    st.session_state.chat_display.append({"role": "user", "content": final_prompt})
    
    with st.spinner("Analyzing data and generating insights..."):
        ai_internal_prompt = final_prompt  
        
        try:
            # --- PPTX GENERATOR ---
            if final_prompt.lower().startswith("deck "):
                csv_filename = final_prompt.split("deck ")[1].strip()
                df = pd.read_csv(csv_filename)
                deck_prompt = f"Analyze this dataset ({csv_filename}):\n{df.head().to_markdown()}\nProvide exactly 3 short bullet points summarizing key insights for a PowerPoint slide. Do NOT output code or markdown asterisks."
                temp_history = [{"role": "system", "content": sys_prompt}, {"role": "user", "content": deck_prompt}]
                deck_resp = client.chat.completions.create(messages=temp_history, model="llama-3.3-70b-versatile")
                slide_text = deck_resp.choices[0].message.content.strip()
                prs = Presentation()
                prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = f"Data Analysis: {csv_filename}"
                slide2 = prs.slides.add_slide(prs.slide_layouts[1])
                slide2.shapes.title.text = "Key Business Insights"
                slide2.shapes.placeholders[1].text_frame.text = slide_text
                pptx_filename = f"Presentation_{csv_filename.replace('.csv', '')}.pptx"
                prs.save(pptx_filename)
                st.success(f"📊 PowerPoint Successfully Generated: {pptx_filename}")
                ai_internal_prompt = f"I generated a PowerPoint presentation for {csv_filename}. Acknowledge this briefly."

            # --- SENTIMENT ANALYSIS ---
            elif final_prompt.lower().startswith("sentiment "):
                csv_filename = final_prompt.split("sentiment ")[1].strip()
                nlp_prompt = f"Write a Python script to read '{csv_filename}'. Find the text column containing the reviews. Use the 'textblob' library (TextBlob(text).sentiment.polarity) to calculate the sentiment of each row. Categorize them as 'Positive' (>0), 'Negative' (<0), or 'Neutral' (==0). Count the categories and plot a pie chart using matplotlib. Save exactly as 'sentiment.png'. Output ONLY valid Python code, no markdown."
                temp_history = st.session_state.conversation_history + [{"role": "user", "content": nlp_prompt}]
                nlp_resp = client.chat.completions.create(messages=temp_history, model="llama-3.3-70b-versatile")
                clean_code = nlp_resp.choices[0].message.content.replace("```python", "").replace("```", "").strip()
                with open("temp_sentiment_script.py", "w") as f: f.write(clean_code)
                subprocess.run(["python", "temp_sentiment_script.py"], check=True)
                st.image("sentiment.png", caption=f"Customer Sentiment for {csv_filename}")
                ai_internal_prompt = "I ran NLP Sentiment Analysis on customer feedback. Briefly explain why this is valuable."

            # --- PDF GENERATOR ---
            elif final_prompt.lower().startswith("report "):
                csv_filename = final_prompt.split("report ")[1].strip()
                df = pd.read_csv(csv_filename)
                report_prompt = f"Write a professional, 3-paragraph executive summary analyzing this dataset ({csv_filename}). Data preview:\n{df.head().to_markdown()}\nIdentify key trends. Do NOT output code or markdown symbols."
                temp_history = [{"role": "system", "content": sys_prompt}, {"role": "user", "content": report_prompt}]
                report_resp = client.chat.completions.create(messages=temp_history, model="llama-3.3-70b-versatile")
                report_text = report_resp.choices[0].message.content.strip()
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Helvetica", style="B", size=16)
                pdf.cell(200, 10, txt=f"Executive Data Report: {csv_filename}", ln=True, align='C')
                pdf.ln(10)
                pdf.set_font("Helvetica", size=12)
                pdf.multi_cell(0, 10, txt=report_text) 
                pdf_filename = f"Executive_Report_{csv_filename.replace('.csv', '')}.pdf"
                pdf.output(pdf_filename)
                st.success(f"📄 PDF Successfully Generated: {pdf_filename}")
                ai_internal_prompt = f"I generated a PDF executive report for {csv_filename}. Acknowledge this briefly."

            # --- MACHINE LEARNING FORECAST ---
            elif final_prompt.lower().startswith("predict "):
                csv_filename = final_prompt.split("predict ")[1].strip()
                ml_prompt = f"Write a Python script to read '{csv_filename}'. 1. Drop missing values. 2. Convert 'Date' to datetime and then to numeric. 3. Train sklearn LinearRegression on Date and Sales. 4. Predict next 30 days. 5. Plot and save as 'forecast.png'. Output ONLY valid Python code."
                temp_history = st.session_state.conversation_history + [{"role": "user", "content": ml_prompt}]
                ml_resp = client.chat.completions.create(messages=temp_history, model="llama-3.3-70b-versatile")
                clean_code = ml_resp.choices[0].message.content.replace("```python", "").replace("```", "").strip()
                with open("temp_predict_script.py", "w") as f: f.write(clean_code)
                subprocess.run(["python", "temp_predict_script.py"], check=True)
                st.image("forecast.png", caption=f"30-Day Forecast for {csv_filename}")
                ai_internal_prompt = "I ran a Machine Learning forecast. Briefly explain how this helps businesses."

            # --- FINANCIAL ANALYST ---
            elif final_prompt.lower().startswith("stock "):
                ticker = final_prompt.split("stock ")[1].strip().upper()
                stock = yf.Ticker(ticker)
                info = stock.info
                price = info.get('currentPrice', 'N/A')
                margins = info.get('profitMargins', 'N/A')
                if margins != 'N/A': margins = f"{margins * 100:.2f}%"
                ai_internal_prompt = f"Live data for {ticker}: Price ${price}, Margins: {margins}. Give a brief analysis."

            # --- VISUALIZER ---
            elif final_prompt.lower().startswith("plot "):
                csv_filename = final_prompt.split("plot ")[1].strip()
                temp_history = st.session_state.conversation_history + [{"role": "user", "content": f"Write Python code to read {csv_filename} and save a chart as 'chart.png'. Only output code."}]
                chart_resp = client.chat.completions.create(messages=temp_history, model="llama-3.3-70b-versatile")
                clean_code = chart_resp.choices[0].message.content.replace("```python", "").replace("```", "").strip()
                with open("temp_chart_script.py", "w") as f: f.write(clean_code)
                subprocess.run(["python", "temp_chart_script.py"], check=True)
                st.image("chart.png")
                ai_internal_prompt = "I generated a chart. Acknowledge this briefly."

            # --- SQL DATABASE ---
            elif final_prompt.lower().startswith("query db "):
                question = final_prompt.split("query db ")[1].strip()
                conn = sqlite3.connect('company_data.db')
                df = pd.read_sql_query("SELECT * FROM employees", conn)
                conn.close()
                st.write("**Database Preview:**")
                st.dataframe(df)
                ai_internal_prompt = f"Database data:\n{df.to_markdown()}\n\nQuestion: {question}. Answer and provide SQL."

            # --- WEB SCRAPER ---
            elif final_prompt.lower().startswith("scrape "):
                url = final_prompt.split("scrape ")[1].strip()
                response = requests.get(url, timeout=10)
                soup = BeautifulSoup(response.content, 'html.parser')
                ai_internal_prompt = f"Scraped text from {url}:\n{soup.get_text(separator=' ', strip=True)[:5000]}\nSummarize this."

            # --- CSV ANALYST ---
            elif "analyze csv" in final_prompt.lower():
                csv_filename = final_prompt.split("analyze csv ")[1].strip()
                df = pd.read_csv(csv_filename)
                st.write("**Data Preview:**")
                st.dataframe(df.head())
                ai_internal_prompt = f"Dataset {csv_filename}:\n{df.head().to_markdown()}\nAnalyze this."

            # 6. Send to AI
            st.session_state.conversation_history.append({"role": "user", "content": ai_internal_prompt})
            
            chat_completion = client.chat.completions.create(
                messages=st.session_state.conversation_history,
                model="llama-3.3-70b-versatile", 
            )
            
            response_text = chat_completion.choices[0].message.content
            
            # File Saving Logic
            if "save it as" in final_prompt.lower() and not any(cmd in final_prompt.lower() for cmd in ["predict ", "plot ", "report ", "sentiment ", "deck "]):
                filename = final_prompt.lower().split("save it as ")[1].split(" ")[0].strip()
                clean_code = response_text.replace("```python", "").replace("```sql", "").replace("```", "").strip()
                with open(filename, "w") as f: f.write(clean_code)
                st.success(f"File saved successfully as {filename}")

            # 7. Print to Screen AND Speak Out Loud (British Accent!)
            with st.chat_message("assistant", avatar="🤖"):
                st.markdown(response_text)
                
                if enable_voice:
                    clean_spoken_text = response_text.replace("*", "").replace("📊", "").replace("📄", "")
                    # Changed tld to 'co.uk' for a professional British accent!
                    tts = gTTS(text=clean_spoken_text, lang='en', tld='co.uk')
                    tts.save("response.mp3")
                    
                    audio_file = open("response.mp3", "rb")
                    audio_bytes = audio_file.read()
                    audio_base64 = base64.b64encode(audio_bytes).decode()
                    audio_html = f"""<audio autoplay><source src="data:audio/mp3;base64,{audio_base64}" type="audio/mp3"></audio>"""
                    st.markdown(audio_html, unsafe_allow_html=True)
                
            st.session_state.chat_display.append({"role": "assistant", "content": response_text})
            st.session_state.conversation_history.append({"role": "assistant", "content": response_text})

        except Exception as e:
            st.error(f"An error occurred: {e}")